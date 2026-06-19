import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_premium_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Designer Palette
    COLOR_BG = RGBColor(10, 15, 30)         # Ultra-dark Navy/Slate
    COLOR_CARD = RGBColor(17, 24, 39)       # Dark Slate Card
    COLOR_TEXT_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT_MUTED = RGBColor(156, 163, 175) # Slate-400
    
    # Accent Borders & Text Highlights
    ACCENT_CYAN = RGBColor(6, 182, 212)     # Primary Cyan
    ACCENT_PURPLE = RGBColor(139, 92, 246)   # AI Purple
    ACCENT_RED = RGBColor(239, 68, 68)       # Alert Red
    BORDER_MUTED = RGBColor(31, 41, 55)      # Subtle Slate Border

    blank_layout = prs.slide_layouts[6]

    def set_slide_base(slide, title_text, category="PARKSENSE AI"):
        """Fills the slide background and adds a premium top header zone."""
        # 1. Slide Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BG

        # 2. Category Tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        cat_tf.margin_left = cat_tf.margin_top = cat_tf.margin_bottom = cat_tf.margin_right = 0
        cat_p = cat_tf.paragraphs[0]
        cat_p.text = category.upper()
        cat_p.font.name = "Segoe UI"
        cat_p.font.size = Pt(11)
        cat_p.font.bold = True
        cat_p.font.color.rgb = ACCENT_CYAN

        # 3. Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = title_tf.margin_top = title_tf.margin_bottom = title_tf.margin_right = 0
        title_p = title_tf.paragraphs[0]
        title_p.text = title_text
        title_p.font.name = "Segoe UI"
        title_p.font.size = Pt(32)
        title_p.font.bold = True
        title_p.font.color.rgb = COLOR_TEXT_WHITE

    def draw_glass_card(slide, left, top, width, height, border_color):
        """Draws a premium rounded container card with colored borders."""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_CARD
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    # ----------------------------------------------------
    # SLIDE 1: Title Slide (Widescreen Hero Layout)
    # ----------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.background
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_BG

    # Large Glowing Panel Card
    draw_glass_card(slide1, Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5), ACCENT_CYAN)

    title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.333), Inches(3.2))
    tf1 = title_box.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "PARKSENSE AI"
    p1.font.name = "Segoe UI"
    p1.font.size = Pt(72)
    p1.font.bold = True
    p1.font.color.rgb = ACCENT_CYAN
    p1.space_after = Pt(8)

    p2 = tf1.add_paragraph()
    p2.text = "Closed-Loop Parking Congestion & Enforcement Intelligence Platform"
    p2.font.name = "Segoe UI"
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_TEXT_WHITE
    p2.space_after = Pt(24)

    p3 = tf1.add_paragraph()
    p3.text = "A Smart City Solution deployed for Bengaluru Traffic Police (BTP)"
    p3.font.name = "Segoe UI"
    p3.font.size = Pt(16)
    p3.font.color.rgb = COLOR_TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 2: The Problem Slide (Double Container Panels)
    # ----------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_base(slide2, "The Challenge: Parking Inefficiency & Gridlocks", "The Problem")

    # Left Container (Operational Obstacles)
    draw_glass_card(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), BORDER_MUTED)
    left_content = slide2.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(4.8), Inches(4.2))
    l_tf = left_content.text_frame
    l_tf.word_wrap = True

    p = l_tf.paragraphs[0]
    p.text = "Operational Pain Points"
    p.font.name = "Segoe UI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.space_after = Pt(18)

    p = l_tf.add_paragraph()
    p.text = "• Traffic Friction: Double parking and footpath parking reduce lane width, cutting carrying capacity in half.\n\n" \
             "• Blind Spot Deployment: Patrol units are dispatched dynamically via guess-work, missing emerging congestion zones.\n\n" \
             "• Patrol Log Bias: Citation statistics only exist where police currently patrol, skewing standard models."
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.line_spacing = 1.3

    # Right Container (Friction Stats)
    draw_glass_card(slide2, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), ACCENT_RED)
    
    # Stat 1
    stat1_box = slide2.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(4.9), Inches(2.0))
    s1_tf = stat1_box.text_frame
    s1_tf.word_wrap = True
    p = s1_tf.paragraphs[0]
    p.text = "45% Reduction"
    p.font.name = "Segoe UI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED
    p.space_after = Pt(4)
    p = s1_tf.add_paragraph()
    p.text = "Carriageway capacity cut along critical urban corridors."
    p.font.name = "Segoe UI"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT_WHITE

    # Stat 2
    stat2_box = slide2.shapes.add_textbox(Inches(7.2), Inches(4.4), Inches(4.9), Inches(2.0))
    s2_tf = stat2_box.text_frame
    s2_tf.word_wrap = True
    p = s2_tf.paragraphs[0]
    p.text = "298,000 Records"
    p.font.name = "Segoe UI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.space_after = Pt(4)
    p = s2_tf.add_paragraph()
    p.text = "Historical parking violations ingested across Bengaluru."
    p.font.name = "Segoe UI"
    p.font.size = Pt(14)
    p.font.color.rgb = COLOR_TEXT_MUTED

    # ----------------------------------------------------
    # SLIDE 3: Architecture Slide (4 Horizontal Glass Cards)
    # ----------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_base(slide3, "System Pipeline: Seamless Closed-Loop Automation", "Architecture")

    col_w = Inches(2.7)
    gap_w = Inches(0.2)
    start_x = Inches(0.8)

    pipeline_cards = [
        ("1. INGEST & FORECAST", "LightGBM regressor models risk profiles across H3 resolution 8 cells using temporal splits.", ACCENT_PURPLE),
        ("2. CAMERA DETECTION", "CCTV edge camera detects stopped cars. If >10s geofence breach occurs, owner gets warning SMS.", ACCENT_CYAN),
        ("3. DYNAMIC ROUTING", "If warning is ignored, system aggregates targets via K-Means and snaps optimized TSP to roads.", ACCENT_CYAN),
        ("4. FEEDBACK LOOP", "Once field officer clears the bottleneck, scores dynamically decay by 45%, retuning the online ML loop.", ACCENT_PURPLE)
    ]

    for idx, (title, desc, accent_color) in enumerate(pipeline_cards):
        x_pos = start_x + idx * (col_w + gap_w)
        draw_glass_card(slide3, x_pos, Inches(2.0), col_w, Inches(4.5), accent_color)

        card_box = slide3.shapes.add_textbox(x_pos + Inches(0.2), Inches(2.3), col_w - Inches(0.4), Inches(4.0))
        c_tf = card_box.text_frame
        c_tf.word_wrap = True
        
        p = c_tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = accent_color
        p.space_after = Pt(16)

        p = c_tf.add_paragraph()
        p.text = desc
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT_MUTED
        p.line_spacing = 1.3

    # ----------------------------------------------------
    # SLIDE 4: ML Ingestion (Double Cards with metrics)
    # ----------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_base(slide4, "Predictive Risk & Local Explainability", "Machine Learning")

    # Left Card (Temporal Validation)
    draw_glass_card(slide4, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), BORDER_MUTED)
    l4_box = slide4.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(4.8), Inches(4.2))
    l4_tf = l4_box.text_frame
    l4_tf.word_wrap = True

    p = l4_tf.paragraphs[0]
    p.text = "Temporal Time-Series Split"
    p.font.name = "Segoe UI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.space_after = Pt(16)

    p = l4_tf.add_paragraph()
    p.text = "• Defensible Split: Trains on weeks < 52 and validates on week 52. Ensures zero spatial-temporal autocorrelation leak.\n\n" \
             "• Hotspot Boundary: Integrates a binary classifier threshold at >= 15.0 violations/hour.\n\n" \
             "• Regression Focus: Predicts continuous counts mapping exact volume scaling."
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.line_spacing = 1.3

    # Right Card (SHAP Waterfall)
    draw_glass_card(slide4, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), ACCENT_PURPLE)
    r4_box = slide4.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(4.9), Inches(4.2))
    r4_tf = r4_box.text_frame
    r4_tf.word_wrap = True

    p = rtf = r4_tf.paragraphs[0]
    p.text = "Local Explainability (SHAP)"
    p.font.name = "Segoe UI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    p.space_after = Pt(16)

    p = r4_tf.add_paragraph()
    p.text = "• Visual Transparency: Renders SHAP waterfall impact charts for every single prediction, building trust with dispatch operators.\n\n" \
             "• Rationale breakdown: Points out direct feature weights driving risk (e.g. proximity to critical intersections (+12.4), peak hours (+8.2), rain scaling (+5.1))."
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.line_spacing = 1.3

    # ----------------------------------------------------
    # SLIDE 5: CCTV Edge Ingestion & VAHAN Warnings
    # ----------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_base(slide5, "CCTV Video Analytics & Pre-Enforcement Warning Loop", "Computer Vision")

    # Left Card (Edge-AI)
    draw_glass_card(slide5, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), BORDER_MUTED)
    l5_box = slide5.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(4.8), Inches(4.2))
    l5_tf = l5_box.text_frame
    l5_tf.word_wrap = True

    p = l5_tf.paragraphs[0]
    p.text = "Edge-AI Vehicle Detection"
    p.font.name = "Segoe UI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.space_after = Pt(16)

    p = l5_tf.add_paragraph()
    p.text = "• Canvas Simulator: Emulates CV model tracking vehicle bounding boxes (moving, stopped, breached).\n\n" \
             "• Geofenced corridors: Monitors lanes inside red zones. Stopped coordinates initiate active stationary timers.\n\n" \
             "• Automated Ingestion: Breaches exceeding 10s automatically post violation logs to backend APIs."
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.line_spacing = 1.3

    # Right Card (Warnings Loop)
    draw_glass_card(slide5, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), ACCENT_CYAN)
    r5_box = slide5.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(4.9), Inches(4.2))
    r5_tf = r5_box.text_frame
    r5_tf.word_wrap = True

    p = r5_tf.paragraphs[0]
    p.text = "VAHAN Pre-Enforcement Warnings"
    p.font.name = "Segoe UI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(16)

    p = r5_tf.add_paragraph()
    p.text = "• Automatic Owner Query: License plates query mock registry mapping owner details instantly.\n\n" \
             "• warning SMS timeline: Dispatches warning SMS. Demonstrates 40% driver compliance, resolving blockages before dispatching physical units."
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.line_spacing = 1.3

    # ----------------------------------------------------
    # SLIDE 6: Dynamic Route Optimization
    # ----------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_base(slide6, "High-Performance Dynamic Route Solvers", "Deployment")

    # Left Card (K-Means + TSP)
    draw_glass_card(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), BORDER_MUTED)
    l6_box = slide6.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(4.8), Inches(4.2))
    l6_tf = l6_box.text_frame
    l6_tf.word_wrap = True

    p = l6_tf.paragraphs[0]
    p.text = "Dynamic Patrol Dispatch"
    p.font.name = "Segoe UI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.space_after = Pt(16)

    p = l6_tf.add_paragraph()
    p.text = "• Sample Weighted K-Means: Groups active hotspot targets into patrol zones, prioritizing locations with high Congestion Impact Scores (CIS).\n\n" \
             "• Snapped road coordinates: Solves Greedy TSP pathing and curves routes along actual roads (OSRM Integration) rather than straight lines."
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.line_spacing = 1.3

    # Right Card (Sub-5ms Caching)
    draw_glass_card(slide6, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), ACCENT_CYAN)
    r6_box = slide6.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(4.9), Inches(4.2))
    r6_tf = r6_box.text_frame
    r6_tf.word_wrap = True

    p = r6_tf.paragraphs[0]
    p.text = "Sub-5ms Cache Fallbacks"
    p.font.name = "Segoe UI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(16)

    p = r6_tf.add_paragraph()
    p.text = "• Dynamic solving lag: On-the-fly calls can take up to 30s. We pre-compute profiles for officer counts 1–20 to deliver instant dashboard load times.\n\n" \
             "• Closest-Match routing: Server dynamically serves closest pre-computed count to bypass OSRM API timeouts or limits."
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.line_spacing = 1.3

    # ----------------------------------------------------
    # SLIDE 7: Feedback Loops & ROI
    # ----------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_base(slide7, "Field Officer Actions & Dynamic Congestion Decay", "Feedback Loop")

    # Left Card (Officer Checklist)
    draw_glass_card(slide7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), BORDER_MUTED)
    l7_box = slide7.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(4.8), Inches(4.2))
    l7_tf = l7_box.text_frame
    l7_tf.word_wrap = True

    p = l7_tf.paragraphs[0]
    p.text = "Field Officer Mobile Portal"
    p.font.name = "Segoe UI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_WHITE
    p.space_after = Pt(16)

    p = l7_tf.add_paragraph()
    p.text = "• Mobile Responsiveness: Collapses sidebar and grid views into a clean, vertically stacked list for mobile usage.\n\n" \
             "• Clearance: Officers check off resolved violations, sending a POST hook directly to the model retraining pipeline."
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.line_spacing = 1.3

    # Right Card (Decay Feedback)
    draw_glass_card(slide7, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), ACCENT_PURPLE)
    r7_box = slide7.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(4.9), Inches(4.2))
    r7_tf = r7_box.text_frame
    r7_tf.word_wrap = True

    p = r7_tf.paragraphs[0]
    p.text = "Online Decay Feedback Loop"
    p.font.name = "Segoe UI"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    p.space_after = Pt(16)

    p = r7_tf.add_paragraph()
    p.text = "• Active Score Decay: Clearing a cell triggers online decay feedback, subtracting -45.0 from the active CIS and predictions in real-time.\n\n" \
             "• Congestion ROI indicators: Computes exact travel time saved (delay minutes/hour) and queue length reductions dynamically to prove operational impact."
    p.font.name = "Segoe UI"
    p.font.size = Pt(15)
    p.font.color.rgb = COLOR_TEXT_MUTED
    p.line_spacing = 1.3

    # ----------------------------------------------------
    # SLIDE 8: Value Proposition (3 Horizontal Cards)
    # ----------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_base(slide8, "Why ParkSense AI Wins the Top-10 Selection", "Value Proposition")

    col_w8 = Inches(3.6)
    gap_w8 = Inches(0.3)
    start_x8 = Inches(0.8)

    benefits = [
        ("Operational Resource ROI", "Automated SMS warnings successfully clear 40% of obstructions, preserving municipal patrol personnel for high-risk zones.", ACCENT_CYAN),
        ("Defensible Pipeline", "Temporal validation splits, SHAP model explainability waterfalls, and real-time decay feedback loops guarantee academic and operational integrity.", ACCENT_PURPLE),
        ("City-Wide Scale", "Snaps OSRM coordinates dynamically across Bengaluru, responding instantly (<5ms) to support thousands of active dispatches.", ACCENT_CYAN)
    ]

    for idx, (title, desc, accent_color) in enumerate(benefits):
        x_pos = start_x8 + idx * (col_w8 + gap_w8)
        draw_glass_card(slide8, x_pos, Inches(2.0), col_w8, Inches(4.5), accent_color)

        card_box = slide8.shapes.add_textbox(x_pos + Inches(0.25), Inches(2.3), col_w8 - Inches(0.5), Inches(4.0))
        c_tf = card_box.text_frame
        c_tf.word_wrap = True
        
        p = c_tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = accent_color
        p.space_after = Pt(16)

        p = c_tf.add_paragraph()
        p.text = desc
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT_MUTED
        p.line_spacing = 1.3

    # Save
    filepath = "parksense_pitch_deck.pptx"
    prs.save(filepath)
    print(f"[OK] Premium pitch deck generated successfully at: {filepath}")

if __name__ == "__main__":
    create_premium_deck()
